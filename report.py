from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def export_ppt(df, output_path="bao_cao_hai_long.pptx"):
    prs = Presentation()

    # =====================
    # SLIDE 1 – TRANG BÌA
    # =====================
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "BÁO CÁO ĐÁNH GIÁ SỰ HÀI LÒNG NGƯỜI BỆNH"
    slide.placeholders[1].text = "BV Đa khoa số 1 tỉnh Lào Cai"

    # =====================
    # SLIDE 2 – TỔNG QUAN
    # =====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. Tổng quan dữ liệu khảo sát"

    total = len(df)
    avg_score = round(df["Do_hai_long"].mean(), 2)

    slide.placeholders[1].text = (
        f"- Tổng số phản hồi: {total}\n"
        f"- Điểm hài lòng trung bình: {avg_score}/5\n"
    )

    # =====================
    # SLIDE 3 – THEO KHOA
    # =====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2. Mức độ hài lòng theo khoa"

    by_khoa = (
        df.groupby("khoa")["Do_hai_long"]
        .mean()
        .sort_values(ascending=False)
    )

    content = ""
    for khoa, diem in by_khoa.items():
        content += f"- {khoa}: {round(diem,2)}/5\n"

    slide.placeholders[1].text = content

    # =====================
    # SLIDE 4 – PHẢN HỒI TIÊU CỰC
    # =====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "3. Phản hồi cần lưu ý"

    bad_df = df[df["Do_hai_long"] <= 2]

    if bad_df.empty:
        slide.placeholders[1].text = "Không ghi nhận phản hồi tiêu cực."
    else:
        text = ""
        for _, r in bad_df.iterrows():
            text += (
                f"- {r['khoa']} | {r['Do_hai_long']} điểm | "
                f"{str(r.get('nguoi_gop_y',''))[:50]}\n"
            )
        slide.placeholders[1].text = text

    # =====================
    # SLIDE 5 – NHẬN XÉT & ĐỀ XUẤT
    # =====================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "4. Nhận xét & đề xuất"

    slide.placeholders[1].text = (
        "- Cần tập trung cải thiện các khoa có điểm trung bình < 3.5\n"
        "- Ưu tiên xử lý phản hồi ≤ 2 điểm trong ngày\n"
        "- Duy trì khảo sát thường xuyên để theo dõi xu hướng\n"
    )

    prs.save(output_path)
    return output_path
